from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
R=Presentation();R.slide_width=Inches(13.333);R.slide_height=Inches(7.5)
BG=RGBColor(248,243,234);INK=RGBColor(35,28,19);ACC=RGBColor(168,67,31);SAGE=RGBColor(113,131,75);MUT=RGBColor(112,98,83);CARD=RGBColor(255,253,248);LINE=RGBColor(221,209,190);GOLD=RGBColor(185,133,35)
def slide():
 s=R.slides.add_slide(R.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=BG
 return s
def text(s,x,y,w,h,value,size=20,color=INK,bold=False,font='Aptos',align=None):
 box=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));tf=box.text_frame;tf.clear();tf.word_wrap=True;p=tf.paragraphs[0];p.text=value;p.font.name=font;p.font.size=Pt(size);p.font.bold=bold;p.font.color.rgb=color
 if align:p.alignment=align
 return box
def rect(s,x,y,w,h,fill=CARD,line=LINE,r=MSO_SHAPE.ROUNDED_RECTANGLE):
 sh=s.shapes.add_shape(r, Inches(x), Inches(y), Inches(w), Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=fill;sh.line.color.rgb=line;return sh
def title(s,kicker,head,sub=''):
 text(s,.65,.38,12,.3,kicker.upper(),10,ACC,True);text(s,.65,.72,12,.65,head,30,INK,True,'Georgia');
 if sub:text(s,.65,1.38,12,.48,sub,15,MUT)
def footer(s,n):text(s,.65,7.12,8,.2,'PashuSetu · SIH26128 · Synthetic prototype data clearly labelled',8,MUT);text(s,12.25,7.08,.4,.25,str(n),9,MUT,True,align=PP_ALIGN.RIGHT)
def card(s,x,y,w,h,num,head,body,color=ACC):
 rect(s,x,y,w,h);rect(s,x+.22,y+.24,.45,.45,color,color,MSO_SHAPE.OVAL);text(s,x+.22,y+.31,.45,.18,str(num),9,RGBColor(255,255,255),True,align=PP_ALIGN.CENTER);text(s,x+.82,y+.22,w-1.05,.35,head,16,INK,True);text(s,x+.82,y+.7,w-1.08,h-.82,body,11,MUT)
#1
s=slide();text(s,.72,.55,5,.3,'SIH26128 · GOVT. OF MAHARASHTRA',11,ACC,True);text(s,.72,1.1,7.8,1.45,'See the outbreak\nbefore it spreads.',42,INK,True,'Georgia');text(s,.75,2.85,6.5,.8,'PashuSetu is the missing last-mile sensor layer for animal-disease surveillance.',19,MUT);rect(s,8.5,.75,4.1,5.65,INK,INK);text(s,9.05,1.25,3,1,'पशुसेतू',34,RGBColor(238,225,191),True,'Georgia');text(s,9.05,2.35,3,.35,'FARMER → SIGNAL → ACTION',11,RGBColor(217,194,145),True);text(s,9.05,3.05,2.8,1.8,'Offline reports\nExplainable triage\nDairy early warning\nFarmer-owned sensors',18,RGBColor(255,255,255),True);text(s,.75,6.45,7,.4,'Preliminary triage, not a diagnosis — consult a veterinarian.',12,ACC,True);footer(s,1)
#2
s=slide();title(s,'The gap','Maharashtra has systems — but no last-mile sensor','Reports arrive late; fragmented records and low connectivity delay containment.');card(s,.65,2.05,3.85,2.25,'01','Farmer signal is missing','NADRES forecasts top-down. Farmers and field workers need a fast offline and voice-first reporting path.');card(s,4.73,2.05,3.85,2.25,'02','Data stays fragmented','Animal identity, symptoms, vaccination, samples and district action do not share one real-time loop.',GOLD);card(s,8.81,2.05,3.85,2.25,'03','Prediction starts too late','Milk-yield and wearable changes can lead obvious clinical signs, but are not connected to surveillance.',SAGE);rect(s,.65,4.75,12,1.2,RGBColor(243,225,215),RGBColor(232,190,172));text(s,1,5.05,11.3,.5,'Our wedge: connect the farmer, the milk and the herd to the officer’s existing workflow — not replace NADRES or Bharat Pashudhan.',18,ACC,True,'Georgia',PP_ALIGN.CENTER);footer(s,2)
#3
s=slide();title(s,'Solution','One secure pipeline, three early signals');
for x,c,h,b in [(0.65,ACC,'Farmer report','Offline app · Marathi/Hindi/English · WhatsApp'),(4.55,GOLD,'Dairy aggregate','Weather-adjusted two-day milk-yield anomaly'),(8.45,SAGE,'Farmer sensor','Private baseline · herd-concurrency detection')]:rect(s,x,2,3.55,1.35);rect(s,x+.25,2.35,.22,.22,c,c,MSO_SHAPE.OVAL);text(s,x+.65,2.2,2.6,.3,h,17,INK,True);text(s,x+.65,2.62,2.55,.45,b,10,MUT)
text(s,.65,3.72,12,.3,'VALIDATE  →  EXPLAINABLE TRIAGE  →  GEO/TIME CLUSTER  →  RIGHT ROLE',14,ACC,True,align=PP_ALIGN.CENTER);rect(s,2.25,4.25,8.8,1.55,INK,INK);text(s,2.7,4.6,7.9,.38,'Farmer sees their animal. Officer sees the district pattern.',22,RGBColor(255,255,255),True,'Georgia',PP_ALIGN.CENTER);text(s,2.7,5.1,7.9,.32,'Raw sensor data never becomes a public dashboard.',11,RGBColor(218,206,190),False,align=PP_ALIGN.CENTER);footer(s,3)
#4
s=slide();title(s,'Defensible mechanisms','Not another report-and-heatmap app');card(s,.65,1.95,5.8,1.35,'A','Covariate-adjusted milk signal','12-week weekday median + robust MAD; removes heat, humidity and rainfall; requires < −2.5σ for two days.',GOLD);card(s,6.75,1.95,5.8,1.35,'B','Correct offline replay','Client UUID, ordered crash-safe queue and idempotent inserts prevent duplicates and loss.',SAGE);card(s,.65,3.55,5.8,1.35,'C','Explainable differential','Top candidates show matched/missed signs, priors, urgency and mandatory safety disclaimer.',ACC);card(s,6.75,3.55,5.8,1.35,'D','The herd reports itself','Per-animal baselines + heat control + multi-animal concurrence create private signals before escalation.',RGBColor(72,110,120));text(s,.65,5.45,12,.45,'Officer confirm/reject closes the safety loop and becomes labelled evidence for future model versions.',17,INK,True,'Georgia',PP_ALIGN.CENTER);footer(s,4)
#5
s=slide();title(s,'Working prototype','One loop from weak signal to containment');steps=[('01','Milk dip','Field verify'),('02','Sensor shift','Private warning'),('03','Farmer report','Offline/voice'),('04','Triage','Reasons shown'),('05','Cluster','District alert'),('06','Case + lab','Contained')]
for i,(n,h,b) in enumerate(steps):
 x=.55+i*2.1;rect(s,x,2.15,1.85,2.1);text(s,x+.18,2.38,.35,.22,n,10,ACC,True);text(s,x+.18,2.85,1.5,.3,h,15,INK,True);text(s,x+.18,3.34,1.5,.45,b,10,MUT)
text(s,.75,4.82,11.8,.4,'Live today',13,ACC,True,align=PP_ALIGN.CENTER);text(s,.75,5.25,11.8,.55,'Next.js PWA · Supabase/PostGIS/RLS · Realtime · deterministic triage · Open-Meteo · production-shaped sensor ingest',16,INK,True,align=PP_ALIGN.CENTER);text(s,.75,6.05,11.8,.4,'₹0 prototype infrastructure · free-tier stack · real hardware swaps into the same JSON contract',13,MUT,align=PP_ALIGN.CENTER);footer(s,5)
#6
s=slide();title(s,'Impact + ask','Pilot one district. Prove earlier verification. Scale through existing systems.');card(s,.65,1.95,3.8,2.2,'1','Pilot','Partner with one Pune dairy network + district veterinary office. Measure time from first signal to field verification.');card(s,4.75,1.95,3.8,2.2,'2','Validate','Track false alerts, sensitivity, farmer completion, officer response and confirmed outcomes.',SAGE);card(s,8.85,1.95,3.8,2.2,'3','Integrate','Import Bharat Pashudhan Tag ID; export NADRS-format reports; expand district by district.',GOLD);rect(s,.65,4.55,12,1.1,RGBColor(243,225,215),RGBColor(232,190,172));text(s,1,4.84,11.3,.48,'Ask: pilot access to anonymised village dairy aggregates, field validation and district workflow feedback.',20,ACC,True,'Georgia',PP_ALIGN.CENTER);text(s,.65,6.12,12,.45,'Safety: decision support only · veterinarian confirmation · DPDP-aware access controls · synthetic prototype data disclosed',11,MUT,align=PP_ALIGN.CENTER);footer(s,6)
R.save('assets/PashuSetu-SIH-Pitch.pptx')
